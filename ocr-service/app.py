import os
import glob
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
import sys

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

# Web framework
from flask import Flask, request, jsonify
from flask_cors import CORS

# Document processing
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from docx import Document
import openpyxl
from pptx import Presentation
import pandas as pd
import fitz  # PyMuPDF for better PDF handling

# Advanced search & AI
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer

# Threading & caching
import threading

app = Flask(__name__)
CORS(app)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max

# Configuration
INDEX_FOLDER = "index_data"
os.makedirs(INDEX_FOLDER, exist_ok=True)

# Initialize AI models
print("Loading AI models...")
try:
    embedding_model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
    print("✓ Embedding model loaded")
except:
    embedding_model = None
    print("⚠ Run: pip install sentence-transformers")

SEARCH_CONFIG = {
    'semantic_weight': 0.6,
    'keyword_weight': 0.4,
    'max_results': 50,
    'enable_ocr': True,
}

# Common English stop words to exclude from keyword/exact-match searches
STOP_WORDS = {
    'a', 'an', 'the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
    'of', 'with', 'by', 'from', 'up', 'about', 'into', 'through', 'during',
    'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had',
    'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might',
    'can', 'that', 'this', 'these', 'those', 'it', 'its', 'as', 'if', 'then',
    'than', 'so', 'yet', 'both', 'either', 'each', 'any', 'all', 'most',
    'more', 'other', 'such', 'no', 'not', 'only', 'own', 'same', 'too',
    'very', 'just', 'because', 'while', 'although', 'however', 'also',
    'he', 'she', 'they', 'we', 'you', 'i', 'me', 'him', 'her', 'us', 'them',
    'my', 'your', 'his', 'our', 'their', 'what', 'which', 'who', 'whom',
    'when', 'where', 'why', 'how', 'over', 'under', 'between', 'after',
    'before', 'above', 'below', 'off', 'out', 'down', 'there', 'here',
}

def filter_stop_words(terms):
    """Remove stop words from a list of terms, keeping meaningful words."""
    filtered = [t for t in terms if t.lower() not in STOP_WORDS and len(t) >= 2]
    # If all words were stop words, return original terms (don't return empty)
    return filtered if filtered else terms

@dataclass
class SearchResult:
    file_id: str
    score: float
    excerpt: str
    matched_terms: List[str] = field(default_factory=list)

class SmartDocumentSearch:
    def __init__(self):
        self.documents_store = {}  # file_id -> Document
        self.embeddings = {}       # file_id -> list of chunk embeddings (numpy arrays)
        # Use token_pattern to keep 1-letter words, remove english stop words to catch specific terms, use ngrams
        self.tfidf_vectorizer = TfidfVectorizer(max_features=10000, token_pattern=r"(?u)\b\w+\b", ngram_range=(1, 2))
        self.tfidf_matrix = None
        self.search_lock = threading.Lock()
        
    @dataclass
    class Document:
        file_id: str
        filepath: str
        content: str
        chunks: List[str]
        metadata: Dict[str, Any]
        
    def extract_text_from_any_file(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        metadata = {
            'filename': os.path.basename(filepath),
            'extracted_method': ''
        }
        ext = filepath.lower().split('.')[-1]
        text = ""
        
        if ext == 'pdf':
            try:
                doc = fitz.open(filepath)
                metadata['extracted_method'] = 'PyMuPDF'
                for page in doc:
                    text += page.get_text()
                    
                if len(text.strip()) < 100 and SEARCH_CONFIG['enable_ocr']:
                    metadata['extracted_method'] = 'OCR (Scanned)'
                    images = convert_from_path(filepath, dpi=200)
                    for img in images:
                        text += pytesseract.image_to_string(img)
            except Exception as e:
                print(f"PDF error: {e}")
                
        elif ext in ['xlsx', 'xls', 'xlsm']:
            try:
                excel_file = pd.ExcelFile(filepath)
                metadata['extracted_method'] = 'Pandas'
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(filepath, sheet_name=sheet_name)
                    text += f"\n=== Sheet: {sheet_name} ===\n"
                    text += df.to_string(index=False)
                    text += " " + " ".join(df.columns.astype(str))
            except Exception as e:
                print(f"Excel error: {e}")
                
        elif ext in ['pptx', 'ppt']:
            try:
                prs = Presentation(filepath)
                metadata['extracted_method'] = 'python-pptx'
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = " | ".join(cell.text for cell in row.cells)
                                text += row_text + "\n"
            except Exception as e:
                print(f"PowerPoint error: {e}")
                
        elif ext == 'docx':
            try:
                doc = Document(filepath)
                metadata['extracted_method'] = 'python-docx'
                text = "\n".join([para.text for para in doc.paragraphs])
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        text += "\n" + row_text
            except Exception as e:
                print(f"Word error: {e}")
                
        elif ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp']:
            try:
                metadata['extracted_method'] = 'OCR'
                img = Image.open(filepath).convert('L')
                if max(img.size) > 2000:
                    ratio = 2000 / max(img.size)
                    new_size = tuple(int(dim * ratio) for dim in img.size)
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                text = pytesseract.image_to_string(img)
            except Exception as e:
                print(f"Image error: {e}")
                
        elif ext in ['txt', 'md', 'json', 'xml', 'html', 'csv']:
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(filepath, 'r', encoding='latin-1') as f:
                    text = f.read()
            metadata['extracted_method'] = 'Direct text'
        else:
            print(f"Unsupported extension: {ext}")
            
        text = ' '.join(text.split())
        return text, metadata
    
    def chunk_document(self, text: str, chunk_size: int = 150, overlap: int = 30) -> List[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if len(chunk.strip()) > 10:  # Avoid tiny empty chunks
                chunks.append(chunk)
        return chunks
    
    def index_document(self, file_id: str, filepath: str) -> bool:
        try:
            content, metadata = self.extract_text_from_any_file(filepath)
            filename = os.path.basename(filepath)
            filename_no_ext = os.path.splitext(filename)[0].lower().replace('-', ' ').replace('_', ' ')

            if not content.strip():
                print(f"⚠ No text extracted from {filepath} — using filename only")
                # For images/unsupported files, use filename as content so TF-IDF works
                content = filename_no_ext

            doc = self.Document(
                file_id=file_id,
                filepath=filepath,
                content=content,
                chunks=self.chunk_document(content),
                metadata={**metadata, 'filename': filename, 'filename_clean': filename_no_ext}
            )
            
            self.documents_store[file_id] = doc
            
            if embedding_model:
                # Encode chunks instead of entire document for much higher accuracy
                chunks_to_encode = doc.chunks[:100] # Limit to 100 chunks per document
                if not chunks_to_encode:
                    chunks_to_encode = [content[:1000]]
                self.embeddings[file_id] = embedding_model.encode(chunks_to_encode)
                
            self.tfidf_matrix = None
            self.save_index()
            print(f"✓ Indexed: {file_id} from {filepath} ({len(doc.chunks)} chunks)")
            return True
            
        except Exception as e:
            print(f"✗ Indexing error for {filepath}: {e}")
            return False
            
    def search_hybrid(self, query: str) -> List[dict]:
        if not self.documents_store:
            return []

        # Filter stop words from query for keyword/exact-match searching
        raw_terms = [t for t in query.lower().split() if len(t) >= 2]
        meaningful_terms = filter_stop_words(raw_terms)
        # Build a cleaned query string using only meaningful terms
        cleaned_query = ' '.join(meaningful_terms)
        effective_query = cleaned_query if cleaned_query else query

        # TF-IDF Search (use cleaned query so stop words don't dilute scoring)
        if self.tfidf_matrix is None:
            documents = [doc.content for doc in self.documents_store.values()]
            if documents:
                self.tfidf_matrix = self.tfidf_vectorizer.fit_transform(documents)
            
        keyword_results = {}
        if self.tfidf_matrix is not None:
            query_vector = self.tfidf_vectorizer.transform([effective_query])
            similarities = cosine_similarity(query_vector, self.tfidf_matrix).flatten()
            for idx, score in enumerate(similarities):
                if score > 0.02:  # lowered from 0.05 for better recall
                    file_id = list(self.documents_store.keys())[idx]
                    keyword_results[file_id] = float(score)
        
        # Semantic Search (Chunk-level) — use original query for better semantic context
        semantic_results = {}
        if self.embeddings and embedding_model:
            query_embedding = embedding_model.encode(effective_query)
            for file_id, chunk_embeddings in self.embeddings.items():
                if len(chunk_embeddings) > 0:
                    # Calculate similarity between query and all chunks of this file
                    similarities = cosine_similarity([query_embedding], chunk_embeddings)[0]
                    best_score = float(np.max(similarities))
                    best_chunk_idx = int(np.argmax(similarities))
                    
                    if best_score > 0.08:  # lowered from 0.15 for better recall
                        semantic_results[file_id] = {
                            'score': best_score,
                            'best_chunk_idx': best_chunk_idx
                        }
                        
        # Direct Substring / Exact Match Search — only on meaningful terms (no stop words)
        query_lower = effective_query.lower()
        query_terms = meaningful_terms  # already filtered
        exact_match_results = {}
        for file_id, doc in self.documents_store.items():
            doc_lower = doc.content.lower()
            exact_score = 0
            
            # High boost for exact full query phrase match
            if query_lower in doc_lower:
                exact_score += 0.8
                
            # Boost for individual partial term matches
            for term in query_terms:
                if term in doc_lower:
                    exact_score += 0.2
                    
            if exact_score > 0:
                exact_match_results[file_id] = min(1.0, exact_score)

        # Filename matching — critical for images and files with little text
        filename_results = {}
        for file_id, doc in self.documents_store.items():
            fname = doc.metadata.get('filename_clean', doc.metadata.get('filename', '')).lower()
            if not fname:
                fname = os.path.splitext(os.path.basename(doc.filepath))[0].lower().replace('-', ' ').replace('_', ' ')
            fname_score = 0
            # Full phrase in filename
            if query_lower in fname:
                fname_score += 0.85
            else:
                # Partial term matches in filename
                for term in query_terms:
                    if term in fname:
                        fname_score += 0.35
            if fname_score > 0:
                filename_results[file_id] = min(1.0, fname_score)

        # Combine all signals
        combined_scores = {}
        all_files = (set(keyword_results.keys()) | set(semantic_results.keys()) |
                     set(exact_match_results.keys()) | set(filename_results.keys()))
        for file_id in all_files:
            kw_score    = keyword_results.get(file_id, 0)
            sem_data    = semantic_results.get(file_id, {'score': 0})
            sem_score   = sem_data['score'] if isinstance(sem_data, dict) else 0
            exact_score = exact_match_results.get(file_id, 0)
            fname_score = filename_results.get(file_id, 0)

            # Filename + exact matches are strongest signals (especially for images)
            combined = (
                SEARCH_CONFIG['keyword_weight'] * kw_score +
                SEARCH_CONFIG['semantic_weight'] * sem_score +
                0.6 * exact_score +
                0.7 * fname_score
            )
            combined_scores[file_id] = min(1.0, combined)
            
        sorted_files = sorted(combined_scores.items(), key=lambda x: x[1], reverse=True)
        
        results = []
        for file_id, score in sorted_files[:SEARCH_CONFIG['max_results']]:
            doc = self.documents_store.get(file_id)
            if doc:
                # Intelligent Excerpt: Use the best semantic chunk if available, otherwise keyword search
                sem_data = semantic_results.get(file_id)
                if sem_data and sem_data['score'] > 0.2 and sem_data['best_chunk_idx'] < len(doc.chunks):
                    excerpt = doc.chunks[sem_data['best_chunk_idx']]
                    if len(excerpt) > 250:
                        excerpt = excerpt[:250] + "..."
                else:
                    excerpt = self.get_excerpt(doc.content, query, 250)
                    
                matched = self.find_matched_terms(doc.content, query)
                results.append({
                    "file_id": file_id,
                    "score": score,
                    "excerpt": excerpt,
                    "matched_terms": matched[:8],
                    "extracted_method": doc.metadata.get('extracted_method', '')
                })
        return results

    def get_excerpt(self, text: str, query: str, length: int = 200) -> str:
        query_lower = query.lower()
        words = text.split()
        best_pos = 0
        best_count = 0
        for i, word in enumerate(words):
            if query_lower in word.lower():
                window = words[max(0, i-50): min(len(words), i+50)]
                count = sum(1 for w in window if query_lower in w.lower())
                if count > best_count:
                    best_count = count
                    best_pos = i
        start = max(0, best_pos - 30)
        end = min(len(words), best_pos + 40)
        excerpt = ' '.join(words[start:end])
        if len(excerpt) > length:
            excerpt = excerpt[:length] + "..."
        return excerpt
    
    def find_matched_terms(self, text: str, query: str) -> List[str]:
        # Filter stop words so only meaningful terms are highlighted
        raw_terms = [t for t in query.lower().split() if len(t) >= 2]
        query_terms = filter_stop_words(raw_terms)
        text_lower = text.lower()
        matched = []
        for term in query_terms:
            if term in text_lower:
                matched.append(term)
        
        # Add the full cleaned phrase if it matches
        cleaned_phrase = ' '.join(query_terms)
        if len(query_terms) > 1 and cleaned_phrase in text_lower:
            matched.insert(0, cleaned_phrase)
            
        return matched
        
    def save_index(self):
        try:
            index_data = {
                'documents': {
                    fid: {
                        'file_id': doc.file_id,
                        'filepath': doc.filepath,
                        'content': doc.content[:2000],
                        'metadata': doc.metadata
                    }
                    for fid, doc in self.documents_store.items()
                }
            }
            with open(os.path.join(INDEX_FOLDER, 'search_index.json'), 'w', encoding='utf-8') as f:
                json.dump(index_data, f, indent=2)
        except Exception as e:
            print(f"Index save error: {e}")
            
    def load_index(self):
        index_path = os.path.join(INDEX_FOLDER, 'search_index.json')
        if os.path.exists(index_path):
            try:
                with open(index_path, 'r', encoding='utf-8') as f:
                    index_data = json.load(f)
                
                # Reconstruct documents
                for fid, data in index_data['documents'].items():
                    self.documents_store[fid] = self.Document(
                        file_id=fid,
                        filepath=data['filepath'],
                        content=data['content'],
                        chunks=[], # Not strictly needed for loaded preview
                        metadata=data['metadata']
                    )
                # Rebuild embeddings from actual files might be slow, normally you'd save embeddings.
                # For this prototype we will rebuild embeddings on load if needed, but it's slow.
                print(f"✓ Loaded index with {len(self.documents_store)} documents")
            except Exception as e:
                print(f"Index load error: {e}")

search_engine = SmartDocumentSearch()
# search_engine.load_index()

@app.route('/api/index', methods=['POST'])
def index_file():
    data = request.json
    if not data or 'filepath' not in data or 'file_id' not in data:
        return jsonify({"success": False, "message": "Missing filepath or file_id"}), 400
        
    file_id = data['file_id']
    filepath = data['filepath']
    
    # Absolute path handling
    if not os.path.isabs(filepath):
        # Assume it's relative to the node backend which passes it
        filepath = os.path.abspath(os.path.join("..", "backend", filepath))
        
    if not os.path.exists(filepath):
        return jsonify({"success": False, "message": f"File not found: {filepath}"}), 404
        
    # Index in background thread to not block the Node.js request
    def run_index():
        search_engine.index_document(file_id, filepath)
        
    threading.Thread(target=run_index).start()
    
    return jsonify({"success": True, "message": "Indexing started"})

@app.route('/api/search', methods=['GET'])
def search_files():
    q = request.args.get('q', '')
    if not q:
        return jsonify({"success": True, "results": []})
        
    results = search_engine.search_hybrid(q)
    return jsonify({
        "success": True,
        "results": results
    })

if __name__ == '__main__':
    app.run(port=5001, debug=True)
